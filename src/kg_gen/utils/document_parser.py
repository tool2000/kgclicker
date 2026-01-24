"""
Document parser module for extracting text from PDF, DOCX, and PPTX files.
"""

from pathlib import Path
from typing import BinaryIO, Optional, Union
import io
import logging

logger = logging.getLogger(__name__)


class DocumentParseError(Exception):
    """Raised when document parsing fails."""

    pass


def extract_text_from_pdf(file: Union[BinaryIO, bytes, Path, str]) -> str:
    """
    Extract text content from a PDF file.

    Args:
        file: File path, bytes, or file-like object

    Returns:
        Extracted text content

    Raises:
        DocumentParseError: If parsing fails
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise DocumentParseError(
            "pymupdf is required for PDF parsing. Install with: pip install pymupdf"
        )

    try:
        if isinstance(file, (str, Path)):
            doc = fitz.open(str(file))
        elif isinstance(file, bytes):
            doc = fitz.open(stream=file, filetype="pdf")
        else:
            # File-like object
            content = file.read()
            doc = fitz.open(stream=content, filetype="pdf")

        text_parts = []
        for page_num, page in enumerate(doc):
            page_text = page.get_text()
            if page_text.strip():
                text_parts.append(f"--- Page {page_num + 1} ---\n{page_text}")

        doc.close()

        if not text_parts:
            logger.warning("PDF appears to be empty or contains only images")
            return ""

        return "\n\n".join(text_parts)

    except Exception as e:
        raise DocumentParseError(f"Failed to parse PDF: {e}") from e


def extract_text_from_docx(file: Union[BinaryIO, bytes, Path, str]) -> str:
    """
    Extract text content from a DOCX file.

    Args:
        file: File path, bytes, or file-like object

    Returns:
        Extracted text content

    Raises:
        DocumentParseError: If parsing fails
    """
    try:
        from docx import Document
    except ImportError:
        raise DocumentParseError(
            "python-docx is required for DOCX parsing. Install with: pip install python-docx"
        )

    try:
        if isinstance(file, (str, Path)):
            doc = Document(str(file))
        elif isinstance(file, bytes):
            doc = Document(io.BytesIO(file))
        else:
            doc = Document(file)

        text_parts = []

        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # Extract tables
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip(" |"):
                    table_text.append(row_text)
            if table_text:
                text_parts.append("\n[Table]\n" + "\n".join(table_text))

        if not text_parts:
            logger.warning("DOCX appears to be empty")
            return ""

        return "\n\n".join(text_parts)

    except Exception as e:
        raise DocumentParseError(f"Failed to parse DOCX: {e}") from e


def extract_text_from_pptx(file: Union[BinaryIO, bytes, Path, str]) -> str:
    """
    Extract text content from a PPTX file.

    Args:
        file: File path, bytes, or file-like object

    Returns:
        Extracted text content

    Raises:
        DocumentParseError: If parsing fails
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise DocumentParseError(
            "python-pptx is required for PPTX parsing. Install with: pip install python-pptx"
        )

    try:
        if isinstance(file, (str, Path)):
            prs = Presentation(str(file))
        elif isinstance(file, bytes):
            prs = Presentation(io.BytesIO(file))
        else:
            prs = Presentation(file)

        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

                # Handle tables in slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells
                        )
                        if row_text.strip(" |"):
                            slide_text.append(row_text)

            if slide_text:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))

        if not text_parts:
            logger.warning("PPTX appears to be empty")
            return ""

        return "\n\n".join(text_parts)

    except Exception as e:
        raise DocumentParseError(f"Failed to parse PPTX: {e}") from e


def extract_text_from_txt(file: Union[BinaryIO, bytes, Path, str]) -> str:
    """
    Extract text content from a TXT file.

    Args:
        file: File path, bytes, or file-like object

    Returns:
        Extracted text content

    Raises:
        DocumentParseError: If parsing fails
    """
    try:
        if isinstance(file, (str, Path)):
            with open(file, "r", encoding="utf-8") as f:
                return f.read()
        elif isinstance(file, bytes):
            return file.decode("utf-8")
        else:
            content = file.read()
            if isinstance(content, bytes):
                return content.decode("utf-8")
            return content
    except UnicodeDecodeError as e:
        raise DocumentParseError(f"Text file must be UTF-8 encoded: {e}") from e
    except Exception as e:
        raise DocumentParseError(f"Failed to read text file: {e}") from e


# Mapping of file extensions to parser functions
PARSERS = {
    ".pdf": extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".doc": extract_text_from_docx,  # Try DOCX parser for .doc (may work for some files)
    ".pptx": extract_text_from_pptx,
    ".ppt": extract_text_from_pptx,  # Try PPTX parser for .ppt (may work for some files)
    ".txt": extract_text_from_txt,
}

SUPPORTED_EXTENSIONS = set(PARSERS.keys())


def get_file_extension(filename: str) -> str:
    """Get lowercase file extension from filename."""
    return Path(filename).suffix.lower()


def is_supported_format(filename: str) -> bool:
    """Check if the file format is supported."""
    return get_file_extension(filename) in SUPPORTED_EXTENSIONS


def extract_text(
    file: Union[BinaryIO, bytes, Path, str],
    filename: Optional[str] = None,
) -> str:
    """
    Extract text from a document file.

    Automatically detects the file type based on extension and uses
    the appropriate parser.

    Args:
        file: File path, bytes, or file-like object
        filename: Original filename (used to detect file type if file is bytes/BinaryIO)

    Returns:
        Extracted text content

    Raises:
        DocumentParseError: If parsing fails or format is unsupported
    """
    # Determine file extension
    if isinstance(file, (str, Path)):
        ext = get_file_extension(str(file))
    elif filename:
        ext = get_file_extension(filename)
    else:
        raise DocumentParseError(
            "filename is required when file is bytes or file-like object"
        )

    if ext not in PARSERS:
        raise DocumentParseError(
            f"Unsupported file format: {ext}. Supported formats: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    parser = PARSERS[ext]
    text = parser(file)

    logger.info(f"Extracted {len(text)} characters from {ext} file")
    return text


def extract_text_from_multiple(
    files: list[tuple[Union[BinaryIO, bytes], str]],
) -> dict[str, str]:
    """
    Extract text from multiple document files.

    Args:
        files: List of (file_content, filename) tuples

    Returns:
        Dictionary mapping filename to extracted text

    Raises:
        DocumentParseError: If any parsing fails
    """
    results = {}
    errors = []

    for file_content, filename in files:
        try:
            text = extract_text(file_content, filename)
            results[filename] = text
        except DocumentParseError as e:
            errors.append(f"{filename}: {e}")

    if errors:
        raise DocumentParseError(
            f"Failed to parse some files:\n" + "\n".join(errors)
        )

    return results
